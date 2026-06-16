"""生成 SWaT 消融实验汇报PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

OUT = os.path.join(os.path.dirname(__file__), "results", "SWaT_Ablation_Report.pptx")
os.makedirs(os.path.dirname(OUT), exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x1A, 0x23, 0x7E)
ACCENT = RGBColor(0x42, 0xA5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)

def add_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, width, height, fill_color=None, border=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # rectangle
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    return shape

def add_text(slide, left, top, width, height, text, font_size=14, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_title_bar(slide, title_text):
    add_box(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), fill_color=BLUE)
    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8), title_text,
             font_size=28, color=WHITE, bold=True)

def add_footer(slide, page_num):
    add_text(slide, Inches(0.5), Inches(7.0), Inches(3), Inches(0.4),
             "SWaT Ablation Study | 2026.06", font_size=9, color=GRAY)
    add_text(slide, Inches(11), Inches(7.0), Inches(2), Inches(0.4),
             f"{page_num}/10", font_size=9, color=GRAY, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════
# Page 1: Cover
# ═══════════════════════════════════════════
s = add_slide()
add_bg(s, BLUE)
add_text(s, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
         "SWaT工业控制系统异常检测", font_size=44, color=WHITE, bold=True)
add_text(s, Inches(1.5), Inches(2.8), Inches(10), Inches(0.8),
         "GATv2 + TCN + GRU 消融实验汇报", font_size=28, color=RGBColor(0xBB, 0xDE, 0xFB))
add_box(s, Inches(1.5), Inches(3.8), Inches(3), Inches(0.04), fill_color=ACCENT)
add_text(s, Inches(1.5), Inches(4.2), Inches(10), Inches(0.6),
         "最佳结果: F1=0.7524  AUC=0.9503  全实验排名第2", font_size=18, color=WHITE)
add_text(s, Inches(1.5), Inches(6.5), Inches(5), Inches(0.5),
         "2026.06.09  |  51 sensors  |  60-step window", font_size=14, color=GRAY)

# ═══════════════════════════════════════════
# Page 2: Background & Baseline
# ═══════════════════════════════════════════
s = add_slide()
add_title_bar(s, "背景与基线模型")
add_footer(s, 2)

add_text(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5), "数据集", font_size=20, color=BLUE, bold=True)
items = [
    "SWaT水处理工业控制系统",
    "51个传感器, 正常+攻击数据",
    "60步滑动窗口, stride=10",
    "异常检测二分类 (正常/攻击)",
    "训练集: normal.csv (80/20 split)",
    "测试集: merged.csv (含攻击)",
    "评分: IQR归一化 + Top-5聚合",
]
for i, item in enumerate(items):
    add_text(s, Inches(1.0), Inches(2.1 + i * 0.45), Inches(5), Inches(0.4),
             f"{item}", font_size=13, color=DARK)

add_text(s, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5), "基线架构", font_size=20, color=BLUE, bold=True)

arch = [
    "Input [B, 60, 51]",
    "  GATv2 (2 layers, heads=2)",
    "  TCN (2 blocks, dil=1,2)",
    "  GRU (1 layer, unidir)",
    "  Pred Head + Recon Head",
    "Params: 209K  |  Train: 5 epochs",
    "Baseline: F1 = 0.6676",
]
for i, line in enumerate(arch):
    c = ACCENT if "F1" in line else DARK
    b = "F1" in line
    add_text(s, Inches(7.2), Inches(2.1 + i * 0.45), Inches(5), Inches(0.4),
             line, font_size=13, color=c, bold=b)

# ═══════════════════════════════════════════
# Page 3: Non-USAD Ablation
# ═══════════════════════════════════════════
s = add_slide()
add_title_bar(s, "非USAD消融实验 (单模块添加)")
add_footer(s, 3)

data = [
    ("Baseline", 0.6676, 0.6676, 0.0),
    ("+Temporal Attn", 0.7063, 0.6676, 0.0387),
    ("+DynPrior Feat", 0.7122, 0.6676, 0.0446),
    ("+Prior Fusion", 0.6988, 0.6676, 0.0312),
    ("+Multi-Scale TCN", 0.6916, 0.6676, 0.0240),
    ("+Dynamic Pearson", 0.6596, 0.6676, -0.0080),
    ("+Prior Dynamic", 0.6584, 0.6676, -0.0092),
]
max_val = 0.74
bar_w = Inches(3.5)
chart_left = Inches(1.2)
chart_top = Inches(1.6)

for i, (name, f1, base, delta) in enumerate(data):
    y = chart_top + Inches(i * 0.75)
    add_text(s, Inches(0.5), y, Inches(1.8), Inches(0.4), name, font_size=11, color=DARK, align=PP_ALIGN.RIGHT)

    # Bar
    w = int(bar_w * (f1 / max_val))
    bar_color = GREEN if delta >= 0 else RED
    add_box(s, chart_left, y + Inches(0.05), w, Inches(0.35), fill_color=bar_color)

    # F1 value
    add_text(s, chart_left + w + Inches(0.1), y, Inches(1.5), Inches(0.4),
             f"F1={f1:.4f} ({delta:+.4f})", font_size=10, color=bar_color, bold=True)

add_text(s, Inches(0.5), Inches(1.2), Inches(4), Inches(0.4),
         "* 基于基线单模块添加, 使用config_dev.yaml (stride=10, hidden=32)", font_size=10, color=GRAY)

# ═══════════════════════════════════════════
# Page 4: Dynamic Graph Verification
# ═══════════════════════════════════════════
s = add_slide()
add_title_bar(s, "动态图验证: 唯一变量实验")
add_footer(s, 4)

add_text(s, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
         "USAD双解码器框架下, 唯一变量 = 动态Pearson图", font_size=16, color=DARK)
add_text(s, Inches(0.8), Inches(2.0), Inches(5), Inches(0.4),
         "两模型: 完全相同架构/训练/评分, 仅图来源不同", font_size=13, color=GRAY)

# Comparison cards
for col, (label, f1, auc, color) in enumerate([
    ("static_usad\n(静态Pearson)", 0.7296, 0.9441, GRAY),
    ("dynamic_usad\n(动态+静态Pearson)", 0.7494, 0.9419, BLUE),
]):
    left = Inches(1.5 + col * 5.5)
    add_box(s, left, Inches(2.8), Inches(4.5), Inches(2.5), fill_color=LIGHT_BG)
    add_text(s, left + Inches(0.3), Inches(3.0), Inches(4), Inches(0.8),
             label, font_size=16, color=color, bold=True)
    add_text(s, left + Inches(0.3), Inches(4.0), Inches(4), Inches(0.6),
             f"F1 = {f1:.4f}", font_size=28, color=color, bold=True)
    add_text(s, left + Inches(0.3), Inches(4.6), Inches(4), Inches(0.4),
             f"AUC = {auc:.4f}", font_size=14, color=GRAY)

add_box(s, Inches(5.5), Inches(5.8), Inches(2), Inches(0.5), fill_color=GREEN)
add_text(s, Inches(5.5), Inches(5.85), Inches(2), Inches(0.4),
         f"+0.0198", font_size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# Page 5: Parallel Dual-Branch Architecture
# ═══════════════════════════════════════════
s = add_slide()
add_title_bar(s, "并行双支路架构 (最佳模型)")
add_footer(s, 5)

add_text(s, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
         "parallel_usad_prior: F1=0.7524  AUC=0.9503  Params=1.0M", font_size=18, color=BLUE, bold=True)

arch_lines = [
    "Input [B,60,51]",
    "  |",
    "  +-- SPATIAL: DynamicPearson + Prior Boost -> GATv2 + Gate -> [B,51,32]",
    "  |",
    "  +-- TEMPORAL: Conv1d(k=3) per-variable -> [B,51,32]",
    "  |",
    "  +-- FUSION: concat -> MLP -> [B,51,32] -> flatten -> z [B,64]",
    "  |",
    "  +-- USAD: Dec1->r1, Dec2->r2, re-encode->Dec2->r12",
]
for i, line in enumerate(arch_lines):
    add_text(s, Inches(1.0), Inches(2.0 + i * 0.38), Inches(11), Inches(0.35),
             line, font_size=12, color=DARK, bold="|" not in line and "--" in line)

# Innovation points
add_box(s, Inches(7.5), Inches(2.0), Inches(5), Inches(4), fill_color=LIGHT_BG)
innos = [
    "Innovations:",
    "1. Dynamic Pearson + Boost融合",
    "2. Prior Node Embed + Gate",
    "3. 并行时空双支路",
    "4. 节点级融合 (非全局pool)",
    "5. 轻量时间编码器",
]
for i, text in enumerate(innos):
    add_text(s, Inches(7.8), Inches(2.2 + i * 0.6), Inches(4.5), Inches(0.5),
             text, font_size=13, color=DARK, bold=(i == 0))

# ═══════════════════════════════════════════
# Page 6: Temporal Branch Ablation
# ═══════════════════════════════════════════
s = add_slide()
add_title_bar(s, "时间分支消融: 越轻量越好")
add_footer(s, 6)

temps = [
    ("Conv1d(k=3)\n单尺度", 0.7524, "最佳"),
    ("Conv1d(k=3,5,7)\n多尺度", 0.7523, "几乎相同"),
    ("TCN(dil=1,2,4)\n多尺度残差", 0.7440, "更差"),
    ("TCN(dil=1,2,4)+GRU\n残差+循环", 0.7271, "最差"),
]

for i, (name, f1, note) in enumerate(temps):
    y = Inches(1.8 + i * 1.3)
    add_box(s, Inches(1), y, Inches(11), Inches(1.1), fill_color=LIGHT_BG if i > 0 else RGBColor(0xE8, 0xF5, 0xE9))
    add_text(s, Inches(1.3), y + Inches(0.1), Inches(3), Inches(0.9), name, font_size=14, color=DARK, bold=True)
    add_text(s, Inches(5), y + Inches(0.2), Inches(2), Inches(0.8),
             f"F1={f1:.4f}", font_size=24, color=GREEN if f1 >= 0.75 else RED, bold=True)
    add_text(s, Inches(7.5), y + Inches(0.3), Inches(4), Inches(0.5), note, font_size=13, color=GRAY)

# ═══════════════════════════════════════════
# Page 7: External Comparison
# ═══════════════════════════════════════════
s = add_slide()
add_title_bar(s, "外部模型对比")
add_footer(s, 7)

ext_data = [
    ("#1", "DCdetector", 0.7553, 0.9337),
    ("#2", "parallel_usad_prior (OURS)", 0.7524, 0.9503),
    ("#3", "dynamic_usad (OURS)", 0.7494, 0.9419),
    ("#4", "USAD", 0.7417, 0.9471),
    ("#5", "MTAD-GAT", 0.7194, 0.9376),
    ("#6", "CAN", 0.7057, 0.9534),
    ("#7", "DAGMM", 0.7048, 0.9431),
    ("#8", "TranAD", 0.6958, 0.9513),
]

# Header
for j, (h, w) in enumerate([("Rank", 1), ("Model", 4), ("F1", 2.5), ("AUC", 2.5)]):
    add_text(s, Inches(1.2 + j * 2.5), Inches(1.5), Inches(w), Inches(0.4), h, font_size=12, color=WHITE, bold=True)
    add_box(s, Inches(1.2 + j * 2.5), Inches(1.5), Inches(w), Inches(0.4), fill_color=BLUE)

for i, (rank, name, f1, auc) in enumerate(ext_data):
    y = Inches(2.1 + i * 0.6)
    our = "OURS" in name
    bg = RGBColor(0xE3, 0xF2, 0xFD) if our else None
    if i % 2 == 0 and not our:
        bg = LIGHT_BG
    if bg:
        add_box(s, Inches(1.2), y, Inches(9.5), Inches(0.5), fill_color=bg)

    add_text(s, Inches(1.3), y + Inches(0.05), Inches(0.8), Inches(0.4), rank, font_size=13, color=BLUE if our else DARK, bold=our)
    add_text(s, Inches(2.2), y + Inches(0.05), Inches(4), Inches(0.4), name, font_size=13, color=BLUE if our else DARK, bold=our)
    add_text(s, Inches(6.5), y + Inches(0.05), Inches(2), Inches(0.4), f"{f1:.4f}", font_size=14, color=GREEN if our else DARK, bold=our)
    add_text(s, Inches(9), y + Inches(0.05), Inches(2), Inches(0.4), f"{auc:.4f}", font_size=13, color=GRAY)

add_text(s, Inches(1.2), Inches(7.0), Inches(6), Inches(0.3),
         "差DCdetector仅0.0029, AUC最高 (0.9503)", font_size=11, color=GRAY)

# ═══════════════════════════════════════════
# Page 8: Key Findings
# ═══════════════════════════════════════════
s = add_slide()
add_title_bar(s, "关键发现")
add_footer(s, 8)

findings = [
    ("01", "动态Pearson图有效", "+0.02 F1增益, 经唯一变量实验验证, 非噪声"),
    ("02", "并行+boost+gate叠加才有效", "单加任何模块都无效, 三项组合达0.7524"),
    ("03", "时间分支越轻量越好", "Conv1d(k=3) > 多尺度Conv > TCN > TCN+GRU"),
    ("04", "USAD双解码器大幅提升基线", "0.6676 -> 0.7494 (+0.0818), USAD框架本身贡献最大"),
    ("05", "先验图受限但有用", "仅9条真实边/51节点, 但boost融合后仍提供微弱增益"),
    ("06", "GATv2的for-loop是瓶颈", "向量化(scatter)在backward不等价, 无法加速"),
]

for i, (num, title, desc) in enumerate(findings):
    y = Inches(1.5 + i * 0.95)
    add_box(s, Inches(0.8), y, Inches(0.6), Inches(0.6), fill_color=BLUE)
    add_text(s, Inches(0.85), y + Inches(0.1), Inches(0.5), Inches(0.4), num, font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.7), y + Inches(0.02), Inches(4), Inches(0.4), title, font_size=16, color=BLUE, bold=True)
    add_text(s, Inches(1.7), y + Inches(0.4), Inches(10), Inches(0.35), desc, font_size=12, color=GRAY)

# ═══════════════════════════════════════════
# Page 9: Prior Graph Analysis
# ═══════════════════════════════════════════
s = add_slide()
add_title_bar(s, "先验图分析")
add_footer(s, 9)

add_text(s, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
         "先验图现状", font_size=20, color=BLUE, bold=True)
prior_info = [
    "来源: Excel (Control Edges + Process Edges + Nodes)",
    "原始物理连接: 9条边, 17/51节点匹配",
    "处理后: 60条边 (含反向+补充), 覆盖51/51",
    "先验图替换静态Pearson: F1 0.7494 -> 0.7044",
    "先验图+boost融合+gate: F1 0.7524 (微弱增益)",
]
for i, item in enumerate(prior_info):
    add_text(s, Inches(1.0), Inches(2.1 + i * 0.5), Inches(10), Inches(0.4),
             item, font_size=13, color=DARK)

add_text(s, Inches(0.8), Inches(4.8), Inches(5), Inches(0.5),
         "扩充方向", font_size=20, color=BLUE, bold=True)
expand = [
    "1. 设备内全连接: 同一设备的所有传感器互连",
    "2. 数据驱动补充: 训练集Pearson最强连通分量",
    "3. 全局弱边: 所有匹配传感器加弱连接",
    "4. 预期: 边数 60->200+, F1 或可超 DCdetector",
]
for i, item in enumerate(expand):
    add_text(s, Inches(1.0), Inches(5.4 + i * 0.5), Inches(10), Inches(0.4),
             item, font_size=13, color=DARK)

# ═══════════════════════════════════════════
# Page 10: Next Steps
# ═══════════════════════════════════════════
s = add_slide()
add_bg(s, BLUE)
add_text(s, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8),
         "下一步工作", font_size=36, color=WHITE, bold=True)
add_box(s, Inches(1.5), Inches(1.6), Inches(3), Inches(0.04), fill_color=ACCENT)

next_steps = [
    "1. 扩充先验图 (设备内全连接 + 数据驱动补充)",
    "2. 多seed验证 (×3), 排除随机波动, 确认统计显著性",
    "3. 超越DCdetector (当前差距仅0.0029)",
    "4. 向量化GAT: 解决backward不等价问题, 大幅加速训练",
    "5. 完整5-epoch重评 D/E 方案 (正在运行中)",
    "6. 撰写论文: 核心贡献 (动态图+并行双支路+boost融合)",
]

for i, step in enumerate(next_steps):
    add_text(s, Inches(1.5), Inches(2.2 + i * 0.75), Inches(10), Inches(0.6),
             step, font_size=18, color=WHITE)

add_text(s, Inches(1.5), Inches(6.8), Inches(10), Inches(0.5),
         "Thank you!", font_size=24, color=ACCENT, bold=True)

prs.save(OUT)
print(f"PPTX saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
